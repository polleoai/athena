"""Local file ingestion for Athena.

Detects file type, extracts text content, copies to the correct raw/ directory,
and returns metadata for wiki page creation.

Tiered support — each tier uses lazy imports with helpful install messages:
  Tier 1 (stdlib):  .md .txt .rst .org .csv .json .yaml .py .js .ts .go .rs .sh .sql
  Tier 2 (pdf):     .pdf — requires pip install athena-brain[pdf]
  Tier 3 (ocr):     .png .jpg .jpeg .gif .webp .svg — requires pip install athena-brain[ocr]
  Tier 4 (office):  .docx .pptx .xlsx — requires pip install athena-brain[office]

Usage:
    from ingest_file import ingest_local_file
    result = ingest_local_file(kb_root, '/path/to/file.pdf')
"""

import os
import re
import shutil
import sys
import time
import json


# ── File type classification ───────────────────────────────────────

TEXT_EXTENSIONS = {
    '.md', '.txt', '.rst', '.org', '.csv', '.json', '.yaml', '.yml',
    '.py', '.js', '.ts', '.go', '.rs', '.sh', '.bash', '.zsh',
    '.sql', '.html', '.xml', '.toml', '.ini', '.cfg', '.conf',
    '.r', '.R', '.jl', '.lua', '.rb', '.php', '.java', '.c', '.cpp',
    '.h', '.hpp', '.swift', '.kt', '.scala', '.hs',
}

PDF_EXTENSIONS = {'.pdf'}

IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.svg', '.bmp', '.tiff'}

OFFICE_EXTENSIONS = {'.docx', '.pptx', '.xlsx', '.doc', '.xls', '.ppt'}

# Where each type goes in raw/ — pulled from bin/config/athena.json so the
# artifact subfolder change in Phase 2 is one file update, not five.
from config import raw_dir_for_source_type as _raw_dir_for
TYPE_MAP = {
    'text':   (_raw_dir_for('webpage'), 'webpage'),
    'code':   (_raw_dir_for('repo'),    'repo'),
    'pdf':    (_raw_dir_for('paper'),   'paper'),
    'image':  (_raw_dir_for('image'),   'image'),
    'office': (_raw_dir_for('paper'),   'paper'),
}

CODE_EXTENSIONS = {
    '.py', '.js', '.ts', '.go', '.rs', '.sh', '.bash', '.zsh',
    '.sql', '.r', '.R', '.jl', '.lua', '.rb', '.php', '.java',
    '.c', '.cpp', '.h', '.hpp', '.swift', '.kt', '.scala', '.hs',
}


def classify_file(filepath):
    """Classify a file by its extension. Returns (category, source_type, raw_dir)."""
    ext = os.path.splitext(filepath)[1].lower()

    if ext in PDF_EXTENSIONS:
        return 'pdf', 'paper', _raw_dir_for('paper')
    elif ext in IMAGE_EXTENSIONS:
        return 'image', 'image', _raw_dir_for('image')
    elif ext in OFFICE_EXTENSIONS:
        return 'office', 'paper', _raw_dir_for('paper')
    elif ext in CODE_EXTENSIONS:
        return 'code', 'repo', _raw_dir_for('repo')
    elif ext in TEXT_EXTENSIONS:
        return 'text', 'webpage', _raw_dir_for('webpage')
    else:
        return None, None, None


# ── Text extraction per tier ───────────────────────────────────────

def _extract_text_file(filepath):
    """Tier 1: Read text-based files directly. No dependencies."""
    try:
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()
    except (IOError, OSError) as e:
        return f'[Error reading file: {e}]'


def _extract_pdf(filepath):
    """Tier 2: Extract text from PDF using pymupdf4llm."""
    try:
        import pymupdf4llm
        return pymupdf4llm.to_markdown(filepath)
    except ImportError:
        return ('[PDF text extraction requires pymupdf4llm]\n'
                'Install with: pip install athena-brain[pdf]\n\n'
                f'File saved to raw/ — install the dependency and re-run kb add to extract text.')
    except Exception as e:
        return f'[PDF extraction error: {e}]'


def _extract_image(filepath):
    """Tier 3: Extract text from image using OCR."""
    try:
        from PIL import Image
        import pytesseract
        img = Image.open(filepath)
        text = pytesseract.image_to_string(img)
        return text.strip() if text.strip() else '[No text detected in image]'
    except ImportError:
        return ('[Image OCR requires pytesseract and Pillow]\n'
                'Install with: pip install athena-brain[ocr]\n'
                'Also install Tesseract OCR: brew install tesseract\n\n'
                f'File saved to raw/ — install dependencies and re-run to extract text.')
    except Exception as e:
        return f'[Image OCR error: {e}]'


def _extract_docx(filepath):
    """Tier 4: Extract text from Word document."""
    try:
        from docx import Document
        doc = Document(filepath)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return '\n\n'.join(paragraphs)
    except ImportError:
        return ('[Word document extraction requires python-docx]\n'
                'Install with: pip install athena-brain[office]\n\n'
                f'File saved to raw/ — install the dependency and re-run to extract text.')
    except Exception as e:
        return f'[DOCX extraction error: {e}]'


def _extract_pptx(filepath):
    """Tier 4: Extract text from PowerPoint."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides_text.append(f'### Slide {i}\n\n' + '\n\n'.join(texts))
        return '\n\n'.join(slides_text)
    except ImportError:
        return ('[PowerPoint extraction requires python-pptx]\n'
                'Install with: pip install athena-brain[office]\n\n'
                f'File saved to raw/ — install the dependency and re-run to extract text.')
    except Exception as e:
        return f'[PPTX extraction error: {e}]'


def _extract_xlsx(filepath):
    """Tier 4: Extract text from Excel spreadsheet."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(filepath, read_only=True, data_only=True)
        sheets_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else '' for c in row]
                if any(c.strip() for c in cells):
                    rows.append(' | '.join(cells))
            if rows:
                # Format as markdown table
                header = rows[0]
                separator = ' | '.join(['---'] * len(rows[0].split(' | ')))
                table = f'### {sheet_name}\n\n| {header} |\n| {separator} |'
                for r in rows[1:]:
                    table += f'\n| {r} |'
                sheets_text.append(table)
        wb.close()
        return '\n\n'.join(sheets_text)
    except ImportError:
        return ('[Excel extraction requires openpyxl]\n'
                'Install with: pip install athena-brain[office]\n\n'
                f'File saved to raw/ — install the dependency and re-run to extract text.')
    except Exception as e:
        return f'[XLSX extraction error: {e}]'


def extract_content(filepath, category):
    """Extract text content based on file category."""
    if category == 'text' or category == 'code':
        return _extract_text_file(filepath)
    elif category == 'pdf':
        return _extract_pdf(filepath)
    elif category == 'image':
        return _extract_image(filepath)
    elif category == 'office':
        ext = os.path.splitext(filepath)[1].lower()
        if ext in ('.docx', '.doc'):
            return _extract_docx(filepath)
        elif ext in ('.pptx', '.ppt'):
            return _extract_pptx(filepath)
        elif ext in ('.xlsx', '.xls'):
            return _extract_xlsx(filepath)
    return '[Unsupported file type]'


# ── Slug generation ────────────────────────────────────────────────

def _make_slug(filepath):
    """Generate a URL-safe slug from a filename."""
    name = os.path.splitext(os.path.basename(filepath))[0]
    slug = re.sub(r'[^\w\s-]', '', name).strip().lower()
    slug = re.sub(r'[-\s]+', '-', slug)
    return slug[:80] if slug else 'file-' + str(int(time.time()))


# ── Main ingestion function ────────────────────────────────────────

def ingest_local_file(kb_root, filepath):
    """Ingest a local file into the knowledge base.

    1. Classifies the file type
    2. Copies/moves to the correct raw/ directory
    3. Extracts text content
    4. Returns metadata dict for wiki page creation

    Args:
        kb_root: vault root path
        filepath: absolute path to the local file

    Returns dict with:
        'success': bool
        'raw_path': relative path to the copied raw file
        'title': extracted or generated title
        'content': extracted text content
        'source_type': paper/repo/webpage/image
        'error': error message if failed
    """
    filepath = os.path.realpath(filepath)

    if not os.path.exists(filepath):
        return {'success': False, 'error': f'File not found: {filepath}'}

    if not os.path.isfile(filepath):
        return {'success': False, 'error': f'Not a file: {filepath}'}

    category, source_type, raw_dir = classify_file(filepath)
    if category is None:
        ext = os.path.splitext(filepath)[1]
        return {'success': False, 'error': f'Unsupported file type: {ext}'}

    # All raw-md output goes through wiki_schema.write_raw_page so file
    # ingest produces structurally identical raws to webpage ingest.
    # Binary side-files (PDF, images, office docs) still get copied
    # alongside the markdown companion, but the .md itself is written
    # by the canonical writer.
    from pathlib import Path as _Path
    sys.path.insert(0, os.path.dirname(__file__))
    from wiki_schema import write_raw_page, make_slug as _ws_make_slug  # type: ignore

    ext = os.path.splitext(filepath)[1].lower()
    title = os.path.splitext(os.path.basename(filepath))[0]
    # Use canonical slug derived from filename (file:// has no URL).
    canonical_slug = _ws_make_slug(title) or 'file-ingest'
    file_url = f"file://{os.path.abspath(filepath)}"

    if category in ('pdf', 'image', 'office'):
        # Binary side-file copy. We still need a stable slug for the
        # binary because its name must match the .md companion. Use the
        # canonical slug for that — keeps both files paired.
        target_dir = os.path.join(kb_root, raw_dir)
        os.makedirs(target_dir, exist_ok=True)
        binary_target = os.path.join(target_dir, canonical_slug + ext)
        if os.path.exists(binary_target):
            return {'success': False, 'error': f'Already exists: {os.path.relpath(binary_target, kb_root)}'}
        shutil.copy2(filepath, binary_target)

        content = extract_content(filepath, category)
        body = (f'- **File:** {canonical_slug}{ext}\n'
                f'- **Type:** {category}\n'
                f'- **Size:** {os.path.getsize(filepath):,} bytes\n'
                f'- **Added:** {time.strftime("%Y-%m-%d")}\n'
                f'- **Content fetched:** yes (local file)\n\n'
                f'## Content\n\n{content}\n')
        try:
            md_path_obj = write_raw_page(
                vault=_Path(kb_root), source_type=source_type, url=file_url,
                title=title, body=body, slug_override=canonical_slug,
                extra_frontmatter={'ingested_via': 'file-ingest', 'category': category},
            )
        except ValueError as e:
            return {'success': False, 'error': f'Schema rejected: {e}'}
        raw_rel = os.path.relpath(str(md_path_obj), kb_root)
    else:
        content = extract_content(filepath, category)
        if ext == '.md':
            heading = re.match(r'^#\s+(.+)$', content, re.MULTILINE)
            if heading:
                title = heading.group(1)
                # Strip the source markdown's H1 — write_raw_page adds its own.
                content = re.sub(r'^#\s+.+\n+', '', content, count=1, flags=re.MULTILINE)
            body = (f'- **File:** {os.path.basename(filepath)}\n'
                    f'- **Added:** {time.strftime("%Y-%m-%d")}\n'
                    f'- **Content fetched:** yes (local file)\n\n'
                    f'## Content\n\n{content}\n')
        else:
            lang = ext.lstrip('.')
            body = (f'- **File:** {os.path.basename(filepath)}\n'
                    f'- **Language:** {lang}\n'
                    f'- **Size:** {os.path.getsize(filepath):,} bytes\n'
                    f'- **Added:** {time.strftime("%Y-%m-%d")}\n'
                    f'- **Content fetched:** yes (local file)\n\n'
                    f'## Content\n\n```{lang}\n{content}\n```\n')
        try:
            md_path_obj = write_raw_page(
                vault=_Path(kb_root), source_type=source_type, url=file_url,
                title=title, body=body, slug_override=canonical_slug,
                extra_frontmatter={'ingested_via': 'file-ingest', 'category': category},
            )
        except ValueError as e:
            return {'success': False, 'error': f'Schema rejected: {e}'}
        raw_rel = os.path.relpath(str(md_path_obj), kb_root)

    return {
        'success': True,
        'raw_path': raw_rel,
        'title': title,
        'content': content[:500],  # first 500 chars for summary
        'source_type': source_type,
        'category': category,
        'slug': slug,
    }
